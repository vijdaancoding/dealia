from __future__ import annotations

import csv as csv_module
import gc
import hashlib
import hmac
import io
import logging
import os
import re
from typing import Any, Iterable

import httpx
from fastapi import Depends, FastAPI, Header, HTTPException
from pydantic import BaseModel

log = logging.getLogger("extractor")
app = FastAPI(title="ticket-attachment-extractor")

MAX_BYTES = int(os.getenv("MAX_BYTES", 40 * 1024 * 1024))
DOWNLOAD_TIMEOUT = float(os.getenv("DOWNLOAD_TIMEOUT", 90))
SAMPLE_ROWS = int(os.getenv("SAMPLE_ROWS", 8))

# Bounds that exist purely to keep peak memory predictable on a small box.
MAX_PDF_PAGES = int(os.getenv("MAX_PDF_PAGES", 100))
MAX_SHEETS = int(os.getenv("MAX_SHEETS", 12))
MAX_COLS = int(os.getenv("MAX_COLS", 40))
MAX_ROWS_SCANNED = int(os.getenv("MAX_ROWS_SCANNED", 20000))

# ---------------------------------------------------------------------------
#  Scanned-PDF OCR via OCR.space
# ---------------------------------------------------------------------------
# OCR is the one thing that genuinely cannot be done cheaply in-process: doing
# it locally means tesseract plus rasterising every page to a bitmap, which is
# what was spiking memory hard enough to get this worker OOM-killed. Handing it
# to an HTTP API moves that entire cost off this container - the request holds
# one file's bytes and waits, nothing more.
#
# Applies to SCANNED PDFs ONLY. Images are never sent anywhere; see HANDLERS.
# Leave OCR_SPACE_API_KEY unset to disable: scanned PDFs then return empty text
# with scanned=true, which is honest and still visible in v_extraction_health.
OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY", "")
OCR_SPACE_URL = os.getenv("OCR_SPACE_URL", "https://api.ocr.space/parse/image")
# The free tier rejects anything over 1 MB. Checking locally first turns a
# guaranteed API rejection into a cheap, clearly-labelled skip.
OCR_SPACE_MAX_BYTES = int(os.getenv("OCR_SPACE_MAX_BYTES", 1024 * 1024))
# Engine 2 handles rotated and lower-quality scans noticeably better than the
# default; both are available on the free tier.
OCR_SPACE_ENGINE = os.getenv("OCR_SPACE_ENGINE", "2")
OCR_SPACE_LANGUAGE = os.getenv("OCR_SPACE_LANGUAGE", "eng")
OCR_SPACE_TIMEOUT = float(os.getenv("OCR_SPACE_TIMEOUT", 120))
# Total pages OCR'd per document. A 200-page scan is not worth 60+ API calls
# against a 500/day free quota; the first pages carry the useful content.
OCR_SPACE_MAX_PAGES = int(os.getenv("OCR_SPACE_MAX_PAGES", 20))
# Pages per request. The free tier only reads the first few pages of any one
# PDF, so batching more than this silently drops the remainder.
OCR_SPACE_MAX_PAGES_PER_REQUEST = int(os.getenv("OCR_SPACE_MAX_PAGES_PER_REQUEST", 3))
# Ceiling on the file we're willing to load into pypdf in order to split it.
OCR_SPACE_SPLIT_MAX_BYTES = int(os.getenv("OCR_SPACE_SPLIT_MAX_BYTES", 20 * 1024 * 1024))

EXTRACTOR_SECRET = os.getenv("EXTRACTOR_SECRET", "")
if not EXTRACTOR_SECRET:
    log.warning("EXTRACTOR_SECRET is not set - /extract is running with NO "
                "authentication. Fine for local-only testing, never for a "
                "publicly reachable deployment.")


def require_secret(x_extractor_secret: str = Header(default="")) -> None:
    # Constant-time compare so response timing can't leak the secret.
    if EXTRACTOR_SECRET and not hmac.compare_digest(x_extractor_secret, EXTRACTOR_SECRET):
        raise HTTPException(status_code=401,
                            detail="invalid or missing X-Extractor-Secret header")


class ExtractRequest(BaseModel):
    extract: bool = True
    url: str = ""
    kind: str = "file"
    name: str = ""
    extension: str = ""
    asset_id: str = ""
    max_chars: int = 40000


def _clean(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\x00", "")
    text = re.sub(r"[ \t\u00a0]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _decode(raw: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


# ---------------------------------------------------------------------------
#  PDF
# ---------------------------------------------------------------------------
def extract_pdf(buf: io.BytesIO, req: ExtractRequest) -> dict[str, Any]:
    import pdfplumber

    pages: list[str] = []
    tables_found = 0
    chars = 0

    with pdfplumber.open(buf) as pdf:
        n_pages = len(pdf.pages)
        for i, page in enumerate(pdf.pages[:MAX_PDF_PAGES], 1):
            try:
                body = page.extract_text() or ""
                # Pipe rows survive chunking far better than the
                # whitespace-aligned text pdfplumber returns by default.
                for tbl in page.extract_tables() or []:
                    tables_found += 1
                    body += "\n\n" + "\n".join(
                        " | ".join((c or "").strip() for c in row) for row in tbl[:40]
                    )
            finally:
                # pdfplumber caches parsed objects per page; without this a
                # long PDF holds every page it has touched in memory at once.
                page.flush_cache()

            if body.strip():
                # The page marker lets a retrieved chunk tell the reader where
                # in a 40-page statement the text came from.
                pages.append(f"[page {i}]\n{body.strip()}")
                chars += len(body)
            if chars > req.max_chars:
                break

    text = _clean("\n\n".join(pages))
    del pages

    # A scanned PDF has essentially NO text layer. Hand it to OCR.space rather
    # than OCRing locally - see the OCR_SPACE_API_KEY comment above.
    scanned = chars < 20 * max(1, min(n_pages, MAX_PDF_PAGES))

    if scanned:
        ocr_text, ocr_status = _ocr_scanned_pdf(buf, req)
        if ocr_text:
            return {"status": "ok", "text": _clean(ocr_text)[: req.max_chars],
                    "pages": n_pages, "scanned": True,
                    "ocr_chars": len(ocr_text), "ocr_status": ocr_status,
                    "ocr_provider": "ocr.space", "needs_ocr": False}
        # No text from OCR.space (not configured, too large, rate-limited, or
        # it simply found nothing). needs_ocr tells the workflow it may still be
        # worth trying its own LLM fallback - that call has to happen in n8n
        # because the model credential lives there, not in this container.
        return {"status": "ok", "text": "", "pages": n_pages, "scanned": True,
                "ocr_chars": 0, "ocr_status": ocr_status, "needs_ocr": True}

    return {"status": "ok", "text": text[: req.max_chars], "pages": n_pages,
            "tables_found": tables_found, "scanned": scanned}



# ---------------------------------------------------------------------------
#  Word
# ---------------------------------------------------------------------------
def _ocr_space_call(data: bytes, req: ExtractRequest, label: str) -> tuple[str, str]:
    """One OCR.space request for one chunk of PDF bytes. Never raises."""
    try:
        with httpx.Client(timeout=OCR_SPACE_TIMEOUT) as client:
            r = client.post(
                OCR_SPACE_URL,
                headers={"apikey": OCR_SPACE_API_KEY},
                files={"file": (req.name or "document.pdf", data, "application/pdf")},
                data={
                    # filetype overrides content-type sniffing, which the API
                    # gets wrong often enough to be worth pinning explicitly.
                    "filetype": "PDF",
                    "OCREngine": OCR_SPACE_ENGINE,
                    "language": OCR_SPACE_LANGUAGE,
                    "isOverlayRequired": "false",
                    "detectOrientation": "true",
                    "scale": "true",
                },
            )
        r.raise_for_status()
        body = r.json()
    except Exception as exc:                      # noqa: BLE001
        log.warning("ocr.space request failed for %s %s: %s", req.asset_id, label, exc)
        return "", f"request_failed: {exc}"[:120]

    if body.get("IsErroredOnProcessing"):
        msg = body.get("ErrorMessage") or body.get("ErrorDetails") or "unknown"
        if isinstance(msg, list):
            msg = "; ".join(str(m) for m in msg)
        log.warning("ocr.space error for %s %s: %s", req.asset_id, label, msg)
        return "", f"api_error: {msg}"[:120]

    # One ParsedResults entry per page; join them in order.
    pages = [str(p.get("ParsedText") or "")
             for p in (body.get("ParsedResults") or [])]
    text = "\n\n".join(t for t in pages if t.strip())
    return text, "ok" if text.strip() else "no_text_found"


def _split_pdf(data: bytes, req: ExtractRequest) -> list[tuple[bytes, str]]:
    """Greedily group pages into sub-PDFs that fit OCR.space's limits.

    Scanned PDFs are page images, so they run large - a 3-page scan routinely
    clears the free tier's 1 MB ceiling. Splitting turns "too big, skipped
    entirely" into several small requests that each succeed.

    Two independent caps apply per request: byte size, and page count (the
    free tier only reads the first few pages of any one PDF, so sending more
    silently loses the rest).
    """
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(io.BytesIO(data))
    total = min(len(reader.pages), OCR_SPACE_MAX_PAGES)

    chunks: list[tuple[bytes, str]] = []
    writer = PdfWriter()
    start = 0
    count = 0

    def flush(end: int) -> None:
        nonlocal writer, count
        if not count:
            return
        out = io.BytesIO()
        writer.write(out)
        chunks.append((out.getvalue(), f"pages {start + 1}-{end}"))
        writer = PdfWriter()
        count = 0

    for i in range(total):
        writer.add_page(reader.pages[i])
        count += 1
        probe = io.BytesIO()
        writer.write(probe)
        size = probe.tell()

        if size > OCR_SPACE_MAX_BYTES and count > 1:
            # Over budget: drop this page back out, emit what fits, and let
            # the page start the next chunk instead.
            writer = PdfWriter()
            for j in range(start, i):
                writer.add_page(reader.pages[j])
            count = i - start
            flush(i)
            start = i
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            count = 1
        elif size > OCR_SPACE_MAX_BYTES:
            # A single page that alone exceeds the limit - nothing to split.
            log.warning("ocr.space: page %d of %s is %d bytes, over limit",
                        i + 1, req.asset_id, size)
            writer = PdfWriter()
            count = 0
            start = i + 1
        elif count >= OCR_SPACE_MAX_PAGES_PER_REQUEST:
            flush(i + 1)
            start = i + 1

    flush(total)
    return chunks


def _ocr_scanned_pdf(buf: io.BytesIO, req: ExtractRequest) -> tuple[str, str]:
    """Send a scanned PDF to OCR.space, splitting it first if needed.

    Never raises: OCR is a best-effort enrichment, and a failure here should
    downgrade the result to 'scanned with no text', not fail the attachment.
    """
    if not OCR_SPACE_API_KEY:
        return "", "not_configured"

    buf.seek(0)
    data = buf.read()

    # Splitting loads the whole document into pypdf and serialises chunks from
    # it, so it is only worth doing for files that comfortably fit in memory.
    # Past this, the split itself would recreate the OOM problem OCR offloading
    # was meant to solve - so decline rather than risk taking the worker down.
    if len(data) > OCR_SPACE_SPLIT_MAX_BYTES:
        del data
        return "", "too_large_for_ocr"

    try:
        chunks = _split_pdf(data, req)
    except Exception as exc:                      # noqa: BLE001
        log.warning("pdf split failed for %s: %s", req.asset_id, exc)
        # Fall back to sending it whole - it may still be under the limit.
        chunks = [(data, "whole")] if len(data) <= OCR_SPACE_MAX_BYTES else []
    finally:
        del data

    if not chunks:
        return "", "too_large_for_ocr"

    texts: list[str] = []
    statuses: list[str] = []
    for chunk_bytes, label in chunks:
        text, status = _ocr_space_call(chunk_bytes, req, label)
        if text:
            texts.append(text)
        statuses.append(status)
        del chunk_bytes

    if texts:
        # Partial success still beats nothing: report ok when any chunk read,
        # but keep the failing statuses visible rather than hiding them.
        failed = [s for s in statuses if s != "ok"]
        return ("\n\n".join(texts),
                "ok" if not failed else f"partial: {failed[0]}"[:120])
    return "", statuses[0] if statuses else "no_text_found"


def extract_docx(buf: io.BytesIO, req: ExtractRequest) -> dict[str, Any]:
    import docx

    d = docx.Document(buf)
    parts: list[str] = []
    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        # Heading style is preserved as markdown so the chunker's
        # heading-awareness works on document text too.
        style = (p.style.name or "").lower()
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "2"
            parts.append(f"{'#' * min(int(level) + 1, 6)} {t}")
        elif style.startswith("list"):
            parts.append(f"- {t}")
        else:
            parts.append(t)

    n_tables = len(d.tables)
    for tbl in d.tables:
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in tbl.rows[:60]]
        if rows:
            parts.append("\n".join(rows))

    n_paras = len(d.paragraphs)
    del d
    return {"status": "ok", "text": _clean("\n\n".join(parts))[: req.max_chars],
            "paragraphs": n_paras, "tables_found": n_tables}


# ---------------------------------------------------------------------------
#  Spreadsheets - structure, not cells
# ---------------------------------------------------------------------------
# A 400-row export embedded cell by cell produces vectors dominated by numbers
# and retrieves for nothing. What people search for is the shape: which sheet,
# which columns, roughly what magnitude.
def _as_number(v: Any) -> float | None:
    if isinstance(v, bool) or v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip().replace(",", "")
    if not s:
        return None
    try:
        return float(s)
    except ValueError:
        return None


def _summarise_rows(rows: Iterable[Any], sheet_name: str) -> dict[str, Any] | None:
    """Single streaming pass: headers, sample, per-column stats. Never holds
    more than SAMPLE_ROWS rows at a time regardless of sheet size."""
    headers: list[str] = []
    sample: list[list[str]] = []
    stats: dict[str, dict[str, float]] = {}
    numeric_n: dict[str, int] = {}
    n_rows = 0

    for idx, row in enumerate(rows):
        vals = list(row)[:MAX_COLS]
        if idx == 0:
            headers = [str(v).strip() if v is not None and str(v).strip() else f"col{i+1}"
                       for i, v in enumerate(vals)]
            continue
        if not headers:
            break
        if all(v is None or not str(v).strip() for v in vals):
            continue

        n_rows += 1
        if len(sample) < SAMPLE_ROWS:
            sample.append([("" if v is None else str(v))[:60] for v in vals])

        for i, v in enumerate(vals[: len(headers)]):
            num = _as_number(v)
            if num is None:
                continue
            key = headers[i]
            s = stats.get(key)
            if s is None:
                stats[key] = {"min": num, "max": num, "sum": num}
            else:
                if num < s["min"]:
                    s["min"] = num
                if num > s["max"]:
                    s["max"] = num
                s["sum"] += num
            numeric_n[key] = numeric_n.get(key, 0) + 1

        if n_rows >= MAX_ROWS_SCANNED:
            break

    if not headers or not n_rows:
        return None

    # A column counts as numeric only if most of its values parsed as numbers -
    # otherwise a text column with a few stray digits would report a bogus sum.
    threshold = max(3, int(0.6 * n_rows))
    col_types = {h: ("number" if numeric_n.get(h, 0) >= threshold else "text")
                 for h in headers}
    numeric_summary = {h: {"min": round(s["min"], 4),
                           "max": round(s["max"], 4),
                           "sum": round(s["sum"], 4)}
                       for h, s in stats.items() if col_types.get(h) == "number"}

    md = ["| " + " | ".join(headers) + " |",
          "| " + " | ".join("---" for _ in headers) + " |"]
    for r in sample:
        cells = (r + [""] * len(headers))[: len(headers)]
        md.append("| " + " | ".join(cells) + " |")

    return {"sheet": str(sheet_name), "rows": n_rows, "cols": len(headers),
            "headers": headers, "column_types": col_types,
            "numeric_summary": numeric_summary, "sample_markdown": "\n".join(md)}


def extract_spreadsheet(buf: io.BytesIO, req: ExtractRequest) -> dict[str, Any]:
    ext = (req.extension or "").lower()
    tables: list[dict[str, Any]] = []

    if ext in ("csv", "tsv"):
        buf.seek(0)
        text = _decode(buf.read())
        reader = csv_module.reader(io.StringIO(text),
                                   delimiter="\t" if ext == "tsv" else ",")
        summary = _summarise_rows(reader, "Sheet1")
        del text
        if summary:
            tables.append(summary)
    else:
        import openpyxl

        # read_only streams rows off disk instead of building a full in-memory
        # object graph; data_only skips formula ASTs and returns cached values.
        wb = openpyxl.load_workbook(buf, read_only=True, data_only=True)
        try:
            for ws in list(wb.worksheets)[:MAX_SHEETS]:
                summary = _summarise_rows(ws.iter_rows(values_only=True), ws.title)
                if summary:
                    tables.append(summary)
        finally:
            wb.close()

    blurbs = [f'Sheet "{t["sheet"]}": {t["rows"]} rows x {t["cols"]} columns. '
              f'Columns: {", ".join(t["headers"])}.' for t in tables]

    return {"status": "ok", "text": _clean("\n".join(blurbs))[: req.max_chars],
            "tables": tables, "sheets": len(tables)}


# ---------------------------------------------------------------------------
#  Slides
# ---------------------------------------------------------------------------
def extract_pptx(buf: io.BytesIO, req: ExtractRequest) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(buf)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        bits = [sh.text.strip() for sh in slide.shapes
                if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if bits or notes:
            parts.append(f"[slide {i}]\n" + "\n".join(bits)
                         + (f"\nSpeaker notes: {notes}" if notes else ""))
    n_slides = len(prs.slides)
    del prs
    return {"status": "ok", "text": _clean("\n\n".join(parts))[: req.max_chars],
            "pages": n_slides}


# ---------------------------------------------------------------------------
#  Plain text
# ---------------------------------------------------------------------------
def extract_text(buf: io.BytesIO, req: ExtractRequest) -> dict[str, Any]:
    buf.seek(0)
    return {"status": "ok", "text": _clean(_decode(buf.read()))[: req.max_chars]}


HANDLERS = {
    "pdf": extract_pdf,
    "document": extract_docx,
    "spreadsheet": extract_spreadsheet,
    "presentation": extract_pptx,
    "text": extract_text,
}


# ---------------------------------------------------------------------------
@app.get("/health")
def health() -> dict[str, Any]:
    return {"ok": True, "handlers": sorted(HANDLERS),
            "pdf_ocr": "ocr.space" if OCR_SPACE_API_KEY else "disabled"}


class _TooBig(Exception):
    pass


def _download(url: str) -> tuple[io.BytesIO, str, int]:
    """Stream to one buffer, hashing as we go.

    The hash is computed from the chunks in flight rather than from the
    finished bytes, so the file never needs to exist twice in memory - the
    previous version's buf.getvalue() silently doubled peak usage for every
    single attachment.
    """
    digest = hashlib.sha256()
    buf = io.BytesIO()
    size = 0
    with httpx.Client(timeout=DOWNLOAD_TIMEOUT, follow_redirects=True) as client:
        with client.stream("GET", url) as resp:
            resp.raise_for_status()
            for chunk in resp.iter_bytes(65536):
                size += len(chunk)
                if size > MAX_BYTES:
                    raise _TooBig
                digest.update(chunk)
                buf.write(chunk)
    buf.seek(0)
    return buf, digest.hexdigest(), size


@app.post("/extract", dependencies=[Depends(require_secret)])
def extract(req: ExtractRequest) -> dict[str, Any]:
    if not req.extract or not req.url:
        return {"status": "skipped", "asset_id": req.asset_id}

    handler = HANDLERS.get(req.kind)
    if handler is None:
        return {"status": "skipped", "asset_id": req.asset_id,
                "error": f"no handler for kind={req.kind}"}

    # Every failure below returns 200. The n8n node has onError set to continue,
    # but a non-2xx would still burn its retries and slow the crawl for a file
    # that is never going to parse.
    try:
        buf, content_hash, size = _download(req.url)
    except _TooBig:
        return {"status": "skipped", "asset_id": req.asset_id,
                "error": "exceeds MAX_BYTES"}
    except Exception as exc:                      # noqa: BLE001
        return {"status": "failed", "asset_id": req.asset_id,
                "error": f"download: {exc}"[:300]}

    if not size:
        return {"status": "failed", "asset_id": req.asset_id, "error": "empty file"}

    try:
        result = handler(buf, req)
    except Exception as exc:                      # noqa: BLE001
        log.exception("extract failed for %s", req.asset_id)
        return {"status": "failed", "asset_id": req.asset_id,
                "error": f"{type(exc).__name__}: {exc}"[:300]}
    finally:
        buf.close()
        del buf
        # One worker serving one ticket at a time means peak RSS is what
        # matters, not throughput. Returning pages promptly keeps the next
        # attachment from starting on top of this one's leftovers.
        gc.collect()

    result.setdefault("status", "ok")
    result["asset_id"] = req.asset_id
    result["size_bytes"] = size
    # The same physical file arrives under different HubSpot file IDs when a
    # pasted screenshot gets re-uploaded on each re-quote; Assemble Ticket uses
    # this hash to collapse those into one indexed asset.
    result["content_hash"] = content_hash
    return result
